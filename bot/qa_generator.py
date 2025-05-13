# bot/qa_generator.py
import os, re, json, textwrap, logging
from pathlib import Path
from typing import List, Dict

try:
    import openai
    openai.api_key = os.environ["OPENAI_API_KEY"]
    _HAS_OPENAI = True
except KeyError:
    _HAS_OPENAI = False

# --- חילוץ טקסט מקבצים ---
def extract_text(filepath: str) -> str:
    fp = Path(filepath)
    if fp.suffix.lower() == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(fp)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    elif fp.suffix.lower() in (".docx", ".doc"):
        import docx
        doc = docx.Document(fp)
        return "\n".join(p.text for p in doc.paragraphs)
    elif fp.suffix.lower() in (".pptx",):
        from pptx import Presentation
        prs = Presentation(fp)
        return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    else:
        return ""

# --- יצירת שאלות ---
def build_qa_from_text(txt: str, n: int = 6) -> List[Dict]:
    """
    מחזיר רשימת dict-ים: {type, question, options, correct}
    חצי multiple (5 תשובות), חצי true_false.
    """
    if _HAS_OPENAI:
        return _qa_via_gpt(txt, n)
    else:
        return _qa_via_placeholder(txt, n)

def _qa_via_gpt(txt: str, n: int):
    prompt = textwrap.dedent(f"""
    צור {n} שאלות בעברית על בסיס הטקסט הבא.
    - חצי מהשאלות מסוג multiple עם 5 אפשרויות (סמני אותיות א.-ה.).
    - חצי True/False.
    החזר JSON בלבד ללא הסברים:
    """) + txt[:5000]
    rsp = openai.ChatCompletion.create(
        model="gpt-4o-mini", messages=[{"role": "user", "content": prompt}]
    )
    # שולף את ה-JSON הראשון בתשובה
    raw = re.search(r"\{.*\}", rsp.choices[0].message.content, re.S).group()
    return json.loads(raw)

def _qa_via_placeholder(txt: str, n: int):
    # אלגוריתם דמה אם אין OpenAI – מייצר שאלה אחת בסיסית
    return [
        {
            "type": "true_false",
            "question": "זהו משפט דוגמה שנוצר כי אין OpenAI API ☺",
            "options": ["נכון", "לא נכון"],
            "correct": "נכון",
        }
    ]
